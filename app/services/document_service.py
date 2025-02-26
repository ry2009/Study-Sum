import os
import PyPDF2
from pptx import Presentation
import docx
import nltk
from nltk.tokenize import sent_tokenize
from nltk.corpus import stopwords
from nltk.cluster.util import cosine_distance
import numpy as np
import networkx as nx
import re

# Download necessary NLTK data
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')

try:
    nltk.data.find('corpora/stopwords')
except LookupError:
    nltk.download('stopwords')

def extract_text_from_pdf(file_path):
    """Extract text from a PDF file"""
    text = ""
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfFileReader(file)
            for page_num in range(pdf_reader.numPages):
                page = pdf_reader.getPage(page_num)
                text += page.extractText() + "\n"
    except Exception as e:
        raise Exception(f"Error extracting text from PDF: {str(e)}")
    return text

def extract_text_from_pptx(file_path):
    """Extract text from a PowerPoint file"""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        raise Exception(f"Error extracting text from PowerPoint: {str(e)}")
    return text

def extract_text_from_docx(file_path):
    """Extract text from a Word document"""
    text = ""
    try:
        doc = docx.Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        raise Exception(f"Error extracting text from Word document: {str(e)}")
    return text

def extract_text_from_txt(file_path):
    """Extract text from a plain text file"""
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    except UnicodeDecodeError:
        # Try another encoding if UTF-8 fails
        with open(file_path, 'r', encoding='latin-1') as file:
            return file.read()
    except Exception as e:
        raise Exception(f"Error extracting text from text file: {str(e)}")

def extract_text_from_document(file_path):
    """Extract text from various document formats with improved encoding handling"""
    file_extension = os.path.splitext(file_path)[1].lower()
    
    try:
        if file_extension == '.pdf':
            return extract_text_from_pdf(file_path)
        elif file_extension in ['.pptx', '.ppt']:
            return extract_text_from_pptx(file_path)
        elif file_extension in ['.docx', '.doc']:
            return extract_text_from_docx(file_path)
        elif file_extension == '.txt':
            return extract_text_from_txt(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_extension}")
    except Exception as e:
        # Try more aggressive text extraction as a fallback
        try:
            with open(file_path, 'rb') as f:
                content = f.read()
                
            # Try multiple encodings
            encodings = ['utf-8', 'latin-1', 'cp1252', 'ascii']
            for encoding in encodings:
                try:
                    text = content.decode(encoding, errors='replace')
                    if text and not text.isspace():
                        return text
                except:
                    continue
                    
            # If still no text, try character by character with ascii filtering
            text = ""
            for char in content:
                if 32 <= char <= 126:  # ASCII printable characters
                    text += chr(char)
            
            if text and not text.isspace():
                return text
            else:
                raise ValueError("Failed to extract text with all methods")
        except:
            raise Exception(f"Error extracting text from document: {str(e)}")

def sentence_similarity(sent1, sent2, stopwords=None):
    """Calculate similarity between two sentences."""
    if stopwords is None:
        stopwords = []
    
    sent1 = [w.lower() for w in sent1]
    sent2 = [w.lower() for w in sent2]
    
    all_words = list(set(sent1 + sent2))
    
    vector1 = [0] * len(all_words)
    vector2 = [0] * len(all_words)
    
    # Build the vectors for both sentences
    for w in sent1:
        if w in stopwords:
            continue
        vector1[all_words.index(w)] += 1
    
    for w in sent2:
        if w in stopwords:
            continue
        vector2[all_words.index(w)] += 1
    
    # Calculate cosine similarity
    return 1 - cosine_distance(vector1, vector2)

def build_similarity_matrix(sentences, stop_words):
    """Build a similarity matrix for all sentences."""
    # Create an empty similarity matrix
    similarity_matrix = np.zeros((len(sentences), len(sentences)))
    
    # Calculate similarity for each pair of sentences
    for i in range(len(sentences)):
        for j in range(len(sentences)):
            if i == j:  # Same sentence, similarity is 1
                similarity_matrix[i][j] = 1
            else:
                similarity_matrix[i][j] = sentence_similarity(
                    sentences[i], sentences[j], stop_words)
    
    return similarity_matrix

def extract_summary_from_text(text, num_sentences=5):
    """Extract summary using enhanced extractive summarization."""
    if not text:
        return "No text available for summarization."
    
    # Tokenize the text into sentences
    sentences = sent_tokenize(text)
    
    # Return the whole text if it's already shorter than requested
    if len(sentences) <= num_sentences:
        return text
    
    # Get stop words
    stop_words = set(stopwords.words('english'))
    
    # Tokenize each sentence into words
    sentence_tokens = [nltk.word_tokenize(sentence) for sentence in sentences]
    
    # Build the similarity matrix
    similarity_matrix = build_similarity_matrix(sentence_tokens, stop_words)
    
    # Create a graph and apply pagerank
    nx_graph = nx.from_numpy_array(similarity_matrix)
    scores = nx.pagerank(nx_graph)
    
    # Sort sentences by score and position
    # Give higher weight to sentences at the beginning and end of the document
    weighted_sentences = []
    for i, (score, sentence) in enumerate(zip(scores, sentences)):
        position_weight = 1.0
        # Increase weight for first few sentences (introduction)
        if i < len(sentences) * 0.1:
            position_weight = 1.5
        # Increase weight for last few sentences (conclusion)
        elif i > len(sentences) * 0.9:
            position_weight = 1.3
        weighted_sentences.append((score * position_weight, i, sentence))
    
    ranked_sentences = sorted(weighted_sentences, reverse=True)
    
    # Get top sentences
    top_sentences = [(idx, sentence) for _, idx, sentence in ranked_sentences[:num_sentences]]
    
    # Sort by their original order in the document
    top_sentences.sort(key=lambda x: x[0])
    ordered_sentences = [sentence for _, sentence in top_sentences]
    
    # Join the sentences
    summary = ' '.join(ordered_sentences)
    
    return summary

def summarize_text(text, max_length=800):
    """Summarize text using advanced extractive summarization."""
    if not text or len(text.strip()) < 100:
        return "Document contains insufficient text for summarization."
    
    try:
        # First, clean the text of any Unicode escape sequences or non-printable chars
        cleaned_text = ""
        for char in text:
            if ord(char) < 128 and (char.isalnum() or char.isspace() or char in '.,;:!?()[]{}-_\'\"'):
                cleaned_text += char
            else:
                cleaned_text += ' '
        
        # If the text has a lot of "/uni..." patterns, it's likely unicode encoded
        if re.search(r'/uni\d{8}', text):
            # Try to clean up unicode escape sequences
            cleaned_text = re.sub(r'/uni\d{8}', ' ', text)
        
        # First check if we have section headers to structure the summary
        lines = cleaned_text.split('\n')
        potential_headers = []
        
        # Look for potential section headers (short lines that end with a colon or are in ALL CAPS)
        for i, line in enumerate(lines):
            stripped = line.strip()
            if (len(stripped) > 0 and len(stripped) < 50 and 
                (stripped.isupper() or stripped.endswith(':') or 
                 stripped.startswith('#') or stripped.startswith('Chapter'))):
                potential_headers.append((i, stripped))
        
        # If we found headers, create a structured summary
        if len(potential_headers) >= 2 and len(potential_headers) <= 10:
            structured_summary = "# Document Outline\n\n"
            
            for i in range(len(potential_headers)):
                header_idx, header = potential_headers[i]
                
                # Add the header
                structured_summary += f"## {header}\n"
                
                # Find the content between this header and the next
                start_idx = header_idx + 1
                end_idx = potential_headers[i+1][0] if i < len(potential_headers) - 1 else len(lines)
                
                # Get the content
                section_content = '\n'.join(lines[start_idx:end_idx])
                
                # Summarize this section (just 1-2 sentences)
                if len(section_content.strip()) > 0:
                    section_summary = extract_summary_from_text(section_content, num_sentences=1)
                    structured_summary += f"{section_summary}\n\n"
            
            if len(structured_summary) < max_length:
                return structured_summary
        
        # Extract key terms for better summarization
        stop_words = set(stopwords.words('english'))
        words = nltk.word_tokenize(cleaned_text.lower())
        filtered_words = [word for word in words if word.isalnum() and word not in stop_words and len(word) > 3]
        
        # Get most common terms
        from collections import Counter
        word_freq = Counter(filtered_words)
        key_terms = [word for word, count in word_freq.most_common(10)]
        
        # Prioritize sentences with key terms
        sentences = sent_tokenize(cleaned_text)
        
        # Score sentences based on key terms
        sentence_scores = []
        for i, sentence in enumerate(sentences):
            score = 0
            # Score based on key terms
            for term in key_terms:
                if term in sentence.lower():
                    score += 1
            
            # Prioritize early and late sentences (intro/conclusion)
            if i < len(sentences) * 0.1:  # First 10% of sentences
                score *= 1.5
            elif i > len(sentences) * 0.9:  # Last 10% of sentences
                score *= 1.3
                
            sentence_scores.append((score, i, sentence))
        
        # Get highest scoring sentences
        top_sentences = sorted(sentence_scores, key=lambda x: x[0], reverse=True)[:max(5, max_length//150)]
        
        # Sort by original order
        top_sentences.sort(key=lambda x: x[1])
        
        # Format the summary
        summary = "# Document Summary\n\n"
        
        # Add key terms section
        if key_terms:
            summary += "## Key Terms\n"
            for term in key_terms[:5]:
                summary += f"• {term.capitalize()}\n"
            summary += "\n"
        
        # Add main content
        summary += "## Main Content\n"
        for _, _, sentence in top_sentences:
            summary += f"{sentence}\n\n"
        
        return summary
        
    except Exception as e:
        return f"Error during summarization: {str(e)}. The document may be in an unsupported format or contains corrupted text."

def process_document(file_path):
    """Process a document file and extract useful information"""
    try:
        # Extract text from the document
        document_text = extract_text_from_document(file_path)
        
        # Try to extract metadata
        metadata = {}
        file_extension = os.path.splitext(file_path)[1].lower()
        
        # For PDFs, try to extract more metadata
        if file_extension == '.pdf':
            try:
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfFileReader(file)
                    info = pdf_reader.getDocumentInfo()
                    if info:
                        metadata['title'] = info.get('/Title', '')
                        metadata['author'] = info.get('/Author', '')
                        metadata['subject'] = info.get('/Subject', '')
                        metadata['creation_date'] = info.get('/CreationDate', '')
            except:
                pass
        
        # Summarize the text
        summary = summarize_text(document_text)
        
        # Extract key terms (simple approach - count word frequency)
        words = nltk.word_tokenize(document_text.lower())
        stop_words = set(stopwords.words('english'))
        filtered_words = [word for word in words if word.isalnum() and word not in stop_words]
        
        # Get word frequency
        word_freq = {}
        for word in filtered_words:
            if len(word) > 3:  # Only include words longer than 3 characters
                word_freq[word] = word_freq.get(word, 0) + 1
        
        # Get top terms
        top_terms = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)[:10]
        
        return {
            "text_sample": document_text[:200] + "..." if len(document_text) > 200 else document_text,
            "summary": summary,
            "file_name": os.path.basename(file_path),
            "metadata": metadata,
            "key_terms": [term for term, _ in top_terms]
        }
    except Exception as e:
        raise Exception(f"Error processing document: {str(e)}") 